"""python-pptx 기반 PPT 생성 서비스."""
import asyncio
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

from app.core.config import settings
from app.core.database import ppt_store, outline_store
from app.models.ppt import PPTJob, PPTStatus
from app.models.outline import SlideOutline, SlideType


# 색상 팔레트
COLOR_TITLE_BG   = RGBColor(0x1E, 0x40, 0xAF)   # 진한 파란색
COLOR_SECTION_BG = RGBColor(0x1E, 0x3A, 0x8A)   # 더 진한 파란색
COLOR_HEADER     = RGBColor(0x1E, 0x40, 0xAF)
COLOR_CLOSING_BG = RGBColor(0x06, 0x5F, 0x46)   # 진한 녹색 (closing 슬라이드)
COLOR_ACCENT     = RGBColor(0x10, 0xB9, 0x81)   # 녹색 포인트
COLOR_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_DARK       = RGBColor(0x1F, 0x29, 0x37)
COLOR_GRAY       = RGBColor(0x6B, 0x72, 0x80)
COLOR_LIGHT_BG   = RGBColor(0xF8, 0xFA, 0xFF)   # 콘텐츠 슬라이드 배경
COLOR_BULLET_DOT = RGBColor(0x3B, 0x82, 0xF6)   # 불릿 강조색


def _add_speaker_notes(slide, notes_text: str) -> None:
    """슬라이드에 발표자 노트를 추가한다."""
    if not notes_text or not notes_text.strip():
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text.strip()


class PPTGeneratorService:
    def _create_presentation(self) -> Presentation:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        return prs

    def _add_title_slide(self, prs: Presentation, title: str, subtitle: str = "", notes: str = "") -> None:
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        # 배경
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_TITLE_BG

        # 장식 바 (하단)
        bar = slide.shapes.add_shape(1, Inches(0), Inches(6.8), Inches(13.33), Inches(0.7))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_ACCENT
        bar.line.fill.background()

        # 제목
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.33), Inches(2.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.color.rgb = COLOR_WHITE

        # 부제목
        if subtitle:
            txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.4), Inches(11.33), Inches(1.2))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.runs[0]
            run2.font.size = Pt(22)
            run2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)

        _add_speaker_notes(slide, notes)

    def _add_section_slide(self, prs: Presentation, title: str, notes: str = "") -> None:
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_SECTION_BG

        # 중앙 구분선
        line = slide.shapes.add_shape(1, Inches(4), Inches(3.55), Inches(5.33), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_ACCENT
        line.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.33), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = COLOR_WHITE

        _add_speaker_notes(slide, notes)

    def _add_content_slide(
        self,
        prs: Presentation,
        title: str,
        bullets: list[str],
        source_texts: list[str] | None = None,
        notes: str = "",
        closing: bool = False,
    ) -> None:
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        # 배경
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_LIGHT_BG if not closing else RGBColor(0xF0, 0xFD, 0xF7)

        header_color = COLOR_CLOSING_BG if closing else COLOR_HEADER

        # 헤더 바
        header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.15))
        header.fill.solid()
        header.fill.fore_color.rgb = header_color
        header.line.fill.background()

        # 헤더 좌측 포인트 바
        accent_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.07), Inches(1.15))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = COLOR_ACCENT
        accent_bar.line.fill.background()

        # 제목
        txTitle = slide.shapes.add_textbox(Inches(0.25), Inches(0.18), Inches(12.5), Inches(0.8))
        tf = txTitle.text_frame
        p = tf.paragraphs[0]
        p.text = title
        run = p.runs[0]
        run.font.size = Pt(26)
        run.font.bold = True
        run.font.color.rgb = COLOR_WHITE

        # 본문 영역 (출처 공간 확보: source 있으면 5.5, 없으면 6.0)
        body_height = Inches(5.5) if source_texts else Inches(6.0)
        txBody = slide.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(12.5), body_height)
        tf = txBody.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

            # 서브불릿 (→ 로 시작)
            is_sub = bullet.strip().startswith("→") or bullet.strip().startswith("  →")
            clean = bullet.strip().lstrip("→").strip()

            if is_sub:
                p.text = f"    › {clean}"
                p.space_before = Pt(2)
                run = p.runs[0]
                run.font.size = Pt(15)
                run.font.color.rgb = COLOR_GRAY
            else:
                p.text = f"• {clean}"
                p.space_before = Pt(8)
                run = p.runs[0]
                run.font.size = Pt(17)
                run.font.bold = False
                run.font.color.rgb = COLOR_DARK

        # 출처 (하단)
        if source_texts:
            source_str = "  |  ".join(source_texts[:4])
            txSrc = slide.shapes.add_textbox(Inches(0.3), Inches(7.05), Inches(12.5), Inches(0.35))
            tf_src = txSrc.text_frame
            p_src = tf_src.paragraphs[0]
            p_src.text = f"📄 {source_str}"
            run_src = p_src.runs[0]
            run_src.font.size = Pt(9)
            run_src.font.color.rgb = COLOR_GRAY

        _add_speaker_notes(slide, notes)

    def _add_references_slide(self, prs: Presentation, sources: list[dict]) -> None:
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xF9, 0xFA, 0xFB)

        header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.15))
        header.fill.solid()
        header.fill.fore_color.rgb = COLOR_ACCENT
        header.line.fill.background()

        txTitle = slide.shapes.add_textbox(Inches(0.25), Inches(0.18), Inches(12), Inches(0.8))
        tf = txTitle.text_frame
        p = tf.paragraphs[0]
        p.text = "참고 자료"
        run = p.runs[0]
        run.font.size = Pt(26)
        run.font.bold = True
        run.font.color.rgb = COLOR_WHITE

        txBody = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.7))
        tf = txBody.text_frame
        tf.word_wrap = True

        seen: set[str] = set()
        ref_list: list[str] = []
        for s in sources:
            name = s.get("source_name") or s.get("name", "")
            if name and name not in seen:
                seen.add(name)
                ref_list.append(name)

        for i, ref in enumerate(ref_list):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"[{i+1}]  {ref}"
            p.space_before = Pt(6)
            run = p.runs[0]
            run.font.size = Pt(17)
            run.font.color.rgb = COLOR_DARK

    async def generate(self, job_id: str) -> None:
        """PPT 생성을 백그라운드에서 실행."""
        job_data = ppt_store.get(job_id)
        if not job_data:
            return

        job = PPTJob.from_dict(job_data)

        try:
            job.status = PPTStatus.generating
            ppt_store.set(job_id, job.to_dict())

            outline_data = outline_store.get(job.outline_id)
            if not outline_data:
                raise ValueError(f"Outline {job.outline_id} not found")

            outline = SlideOutline.from_dict(outline_data)
            prs = await asyncio.get_event_loop().run_in_executor(
                None, self._build_pptx, outline
            )

            output_dir = Path("./data/ppts")
            output_dir.mkdir(parents=True, exist_ok=True)
            output_path = output_dir / f"{job_id}.pptx"

            await asyncio.get_event_loop().run_in_executor(
                None, prs.save, str(output_path)
            )

            job.status = PPTStatus.done
            job.file_path = str(output_path)
            ppt_store.set(job_id, job.to_dict())

        except Exception as e:
            job.status = PPTStatus.error
            job.error_message = str(e)
            ppt_store.set(job_id, job.to_dict())
            raise

    def _build_pptx(self, outline: SlideOutline) -> Presentation:
        prs = self._create_presentation()
        all_sources: list[dict] = []

        for slide in outline.slides:
            slide_sources = [
                {"source_name": s.source_name, "page": s.page}
                for s in slide.sources
            ]
            all_sources.extend(slide_sources)
            source_texts = sorted({s.source_name for s in slide.sources if s.source_name})
            notes = slide.notes or ""

            if slide.type == SlideType.title:
                subtitle = slide.bullets[0] if slide.bullets else ""
                self._add_title_slide(prs, slide.title, subtitle, notes)

            elif slide.type == SlideType.section:
                self._add_section_slide(prs, slide.title, notes)

            elif slide.type == SlideType.content:
                self._add_content_slide(prs, slide.title, slide.bullets, source_texts or None, notes)

            elif slide.type == SlideType.closing:
                self._add_content_slide(prs, slide.title, slide.bullets, source_texts or None, notes, closing=True)

            # references 타입은 아래에서 별도 생성

        # 참고 자료 슬라이드 항상 마지막
        self._add_references_slide(prs, all_sources)
        return prs


ppt_generator_service = PPTGeneratorService()
