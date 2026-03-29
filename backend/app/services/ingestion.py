"""파일/URL 파싱 및 청킹 서비스."""
import asyncio
from pathlib import Path
from typing import AsyncGenerator

import fitz  # PyMuPDF
import httpx
from bs4 import BeautifulSoup
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation

from app.core.config import settings
from app.core.exceptions import UnsupportedFileTypeError
from app.models.source import Source, SourceType


class ParsedChunk:
    def __init__(self, text: str, source_id: str, page: int | None = None, metadata: dict | None = None):
        self.text = text
        self.source_id = source_id
        self.page = page
        self.metadata = metadata or {}


class IngestionService:
    def __init__(self):
        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=settings.chunk_size,
            chunk_overlap=settings.chunk_overlap,
            separators=["\n\n", "\n", ".", " ", ""],
        )

    async def parse_and_chunk(self, source: Source) -> list[ParsedChunk]:
        """소스 타입에 따라 파싱 후 청크로 분리."""
        if source.type == SourceType.pdf:
            return await self._parse_pdf(source)
        elif source.type == SourceType.pptx:
            return await self._parse_pptx(source)
        elif source.type == SourceType.url:
            return await self._parse_url(source)
        else:
            raise UnsupportedFileTypeError(f"Unsupported source type: {source.type}", source.id)

    async def _parse_pdf(self, source: Source) -> list[ParsedChunk]:
        path = Path(source.original_path)
        chunks: list[ParsedChunk] = []

        # CPU bound 작업은 스레드풀에서 실행
        def _extract():
            doc = fitz.open(str(path))
            result = []
            for page_num, page in enumerate(doc, start=1):
                text = page.get_text("text").strip()
                if text:
                    result.append((text, page_num))
            doc.close()
            return result

        pages = await asyncio.get_event_loop().run_in_executor(None, _extract)

        for text, page_num in pages:
            for chunk_text in self.splitter.split_text(text):
                if chunk_text.strip():
                    chunks.append(ParsedChunk(
                        text=chunk_text.strip(),
                        source_id=source.id,
                        page=page_num,
                        metadata={"source_name": source.name, "page": page_num},
                    ))

        return chunks

    async def _parse_pptx(self, source: Source) -> list[ParsedChunk]:
        path = Path(source.original_path)
        chunks: list[ParsedChunk] = []

        def _extract():
            prs = Presentation(str(path))
            result = []
            for slide_num, slide in enumerate(prs.slides, start=1):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                texts.append(text)
                if texts:
                    result.append(("\n".join(texts), slide_num))
            return result

        slides = await asyncio.get_event_loop().run_in_executor(None, _extract)

        for text, slide_num in slides:
            for chunk_text in self.splitter.split_text(text):
                if chunk_text.strip():
                    chunks.append(ParsedChunk(
                        text=chunk_text.strip(),
                        source_id=source.id,
                        page=slide_num,
                        metadata={"source_name": source.name, "slide": slide_num},
                    ))

        return chunks

    async def _parse_url(self, source: Source) -> list[ParsedChunk]:
        headers = {
            "User-Agent": (
                "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
                "AppleWebKit/537.36 (KHTML, like Gecko) "
                "Chrome/120.0.0.0 Safari/537.36"
            )
        }
        async with httpx.AsyncClient(timeout=30.0, follow_redirects=True, headers=headers) as client:
            response = await client.get(source.original_path)
            response.raise_for_status()

        soup = BeautifulSoup(response.text, "html.parser")

        # 불필요한 태그 제거
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()

        text = soup.get_text(separator="\n", strip=True)
        chunks: list[ParsedChunk] = []

        for chunk_text in self.splitter.split_text(text):
            if chunk_text.strip():
                chunks.append(ParsedChunk(
                    text=chunk_text.strip(),
                    source_id=source.id,
                    page=None,
                    metadata={"source_name": source.name, "url": source.original_path},
                ))

        return chunks


ingestion_service = IngestionService()
